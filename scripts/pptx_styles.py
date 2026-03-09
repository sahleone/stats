"""
pptx_styles.py — Shared python-pptx formatting helpers.

Provides a StatsPresentation class that pre-configures slide layouts,
fonts, and colour schemes for consistent lecture slides across all topics.

Usage:
    from scripts.pptx_styles import StatsPresentation

    prs = StatsPresentation("One-Sample Z-Test")
    prs.add_title_slide("One-Sample Z-Test", "Undergraduate Statistics")
    prs.add_content_slide("When Do We Use This Test?", [
        "We have a single sample from a population",
        "We know the population standard deviation σ",
        "We want to test a claim about the population mean μ",
    ])
    prs.add_formula_slide("Test Statistic", "z = (x̄ − μ₀) / (σ / √n)", [
        "x̄ = sample mean",
        "μ₀ = hypothesized population mean",
        "σ = known population standard deviation",
        "n = sample size",
    ])
    prs.save("one_sample_z_test.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Colours
DARK_BLUE = RGBColor(47, 84, 150)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_GREY = RGBColor(242, 242, 242)
MEDIUM_GREY = RGBColor(128, 128, 128)


class StatsPresentation:
    """Pre-configured presentation builder for stats lecture slides."""

    def __init__(self, topic_title="Hypothesis Test"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.topic_title = topic_title

    def save(self, filepath):
        self.prs.save(filepath)

    # ---- Slide Builders ----

    def add_title_slide(self, title, subtitle="Undergraduate Statistics"):
        """Slide 1: Big title + subtitle on a dark blue background."""
        slide = self._add_blank_slide()
        # Dark blue background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        # Title
        txbox = slide.shapes.add_textbox(
            Inches(1), Inches(2.2), Inches(11.333), Inches(2)
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(180, 199, 231)
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(12)

        return slide

    def add_content_slide(self, title, bullets, notes=None):
        """Standard slide with title and bullet points."""
        slide = self._add_blank_slide()
        self._add_slide_title(slide, title)

        body_top = Inches(1.6)
        txbox = slide.shapes.add_textbox(
            Inches(0.8), body_top, Inches(11.5), Inches(5.2)
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(22)
            p.font.color.rgb = BLACK
            p.space_before = Pt(8)
            p.space_after = Pt(4)
            p.level = 0

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_formula_slide(self, title, formula, definitions=None):
        """Slide with a formula in a grey box and optional symbol definitions."""
        slide = self._add_blank_slide()
        self._add_slide_title(slide, title)

        # Formula box
        box_left = Inches(0.8)
        box_top = Inches(1.8)
        box_w = Inches(11.5)
        box_h = Inches(1.5)
        shape = slide.shapes.add_shape(
            1, box_left, box_top, box_w, box_h  # 1 = rectangle
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GREY
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = formula
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER

        # Definitions below
        if definitions:
            def_top = Inches(3.6)
            txbox = slide.shapes.add_textbox(
                Inches(1.2), def_top, Inches(10.5), Inches(3.5)
            )
            tf2 = txbox.text_frame
            tf2.word_wrap = True

            header_p = tf2.paragraphs[0]
            header_p.text = "where:"
            header_p.font.size = Pt(20)
            header_p.font.italic = True
            header_p.font.color.rgb = MEDIUM_GREY
            header_p.space_after = Pt(6)

            for defn in definitions:
                p = tf2.add_paragraph()
                p.text = f"  {defn}"
                p.font.size = Pt(20)
                p.font.color.rgb = BLACK
                p.space_before = Pt(4)

        return slide

    def add_two_column_slide(self, title, left_title, left_items,
                              right_title, right_items):
        """Slide with two columns (e.g. two-tailed vs one-tailed)."""
        slide = self._add_blank_slide()
        self._add_slide_title(slide, title)

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2)
        )
        self._fill_column(left_box, left_title, left_items)

        # Vertical divider
        slide.shapes.add_shape(
            1, Inches(6.5), Inches(1.8), Inches(0.02), Inches(4.5)
        ).fill.solid()

        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.2)
        )
        self._fill_column(right_box, right_title, right_items)

        return slide

    def add_worked_example_slide(self, title, content_lines):
        """Slide for a worked example step (smaller font, more lines)."""
        slide = self._add_blank_slide()
        self._add_slide_title(slide, title)

        txbox = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.4)
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(content_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Check for bold markers
            if line.startswith("**") and line.endswith("**"):
                p.text = line.strip("*")
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = DARK_BLUE
            else:
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = BLACK
            p.space_before = Pt(4)
            p.space_after = Pt(2)

        return slide

    def add_summary_slide(self, key_points, see_also=None):
        """Final summary slide with key takeaways and optional See Also."""
        slide = self._add_blank_slide()
        self._add_slide_title(slide, "Summary & Key Takeaways")

        txbox = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.5)
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        for i, point in enumerate(key_points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"✓  {point}"
            p.font.size = Pt(22)
            p.font.color.rgb = BLACK
            p.space_before = Pt(6)

        if see_also:
            sa_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8)
            )
            tf2 = sa_box.text_frame
            tf2.word_wrap = True
            header = tf2.paragraphs[0]
            header.text = "See Also:"
            header.font.size = Pt(20)
            header.font.bold = True
            header.font.color.rgb = DARK_BLUE
            header.space_after = Pt(4)

            for name, desc in see_also:
                p = tf2.add_paragraph()
                p.text = f"→ {name} — {desc}"
                p.font.size = Pt(18)
                p.font.color.rgb = MEDIUM_GREY
                p.space_before = Pt(2)

        return slide

    # ---- Internal Helpers ----

    def _add_blank_slide(self):
        layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(layout)

    def _add_slide_title(self, slide, title):
        txbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1.1)
        )
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.LEFT

        # Underline bar
        slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.35), Inches(12), Inches(0.04)
        ).fill.solid()

    def _fill_column(self, textbox, title, items):
        tf = textbox.text_frame
        tf.word_wrap = True

        header = tf.paragraphs[0]
        header.text = title
        header.font.size = Pt(24)
        header.font.bold = True
        header.font.color.rgb = DARK_BLUE
        header.space_after = Pt(8)

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = BLACK
            p.space_before = Pt(4)
